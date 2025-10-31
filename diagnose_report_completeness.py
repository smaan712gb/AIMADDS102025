"""
Comprehensive Report Completeness Diagnostic
Analyzes generated reports to identify gaps and improvement areas
"""

import openpyxl
import json
from pathlib import Path
from datetime import datetime

def analyze_excel_report(filepath):
    """Analyze Excel report structure and content"""
    print("\n" + "="*80)
    print("EXCEL REPORT ANALYSIS")
    print("="*80)
    
    wb = openpyxl.load_workbook(filepath)
    
    analysis = {
        'total_sheets': len(wb.worksheets),
        'sheets': {}
    }
    
    for sheet in wb.worksheets:
        sheet_name = sheet.title
        max_row = sheet.max_row
        max_col = sheet.max_column
        
        # Sample first few rows to check for data
        sample_data = []
        for row in sheet.iter_rows(min_row=1, max_row=min(10, max_row), values_only=True):
            sample_data.append([str(cell) if cell is not None else '' for cell in row])
        
        # Check for empty sheets
        has_data = any(any(cell for cell in row if cell and cell != 'None') for row in sample_data)
        
        analysis['sheets'][sheet_name] = {
            'rows': max_row,
            'columns': max_col,
            'has_data': has_data,
            'sample': sample_data[:3] if has_data else []
        }
        
        print(f"\n{sheet_name}:")
        print(f"  Rows: {max_row}, Columns: {max_col}")
        print(f"  Has Data: {has_data}")
        if has_data and sample_data:
            print(f"  Headers: {sample_data[0][:5]}")  # First 5 headers
    
    return analysis

def analyze_ppt_report(filepath):
    """Analyze PowerPoint report structure"""
    print("\n" + "="*80)
    print("POWERPOINT REPORT ANALYSIS")
    print("="*80)
    
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        
        analysis = {
            'total_slides': len(prs.slides),
            'slides': []
        }
        
        for i, slide in enumerate(prs.slides, 1):
            slide_info = {
                'slide_number': i,
                'title': '',
                'shapes_count': len(slide.shapes),
                'has_text': False,
                'has_table': False,
                'has_chart': False
            }
            
            # Get slide title
            if slide.shapes.title:
                slide_info['title'] = slide.shapes.title.text
            
            # Check for content types
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_info['has_text'] = True
                if shape.has_table:
                    slide_info['has_table'] = True
                if shape.has_chart:
                    slide_info['has_chart'] = True
            
            analysis['slides'].append(slide_info)
            print(f"\nSlide {i}: {slide_info['title']}")
            print(f"  Shapes: {slide_info['shapes_count']}, Text: {slide_info['has_text']}, Table: {slide_info['has_table']}, Chart: {slide_info['has_chart']}")
        
        return analysis
    except ImportError:
        print("python-pptx not available, skipping PPT analysis")
        return {'error': 'python-pptx not installed'}

def analyze_pdf_report(filepath):
    """Analyze PDF report structure"""
    print("\n" + "="*80)
    print("PDF REPORT ANALYSIS")
    print("="*80)
    
    try:
        import PyPDF2
        with open(filepath, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            analysis = {
                'total_pages': len(pdf_reader.pages),
                'pages': []
            }
            
            for i, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                word_count = len(text.split())
                
                # Extract section headers (lines with less than 60 chars)
                lines = text.split('\n')
                headers = [line.strip() for line in lines if line.strip() and len(line.strip()) < 60 and len(line.strip()) > 5][:5]
                
                analysis['pages'].append({
                    'page_number': i,
                    'word_count': word_count,
                    'headers': headers
                })
                
                print(f"\nPage {i}:")
                print(f"  Words: {word_count}")
                if headers:
                    print(f"  Headers: {headers[:3]}")
            
            return analysis
    except ImportError:
        print("PyPDF2 not available, skipping PDF analysis")
        return {'error': 'PyPDF2 not installed'}

def identify_critical_gaps():
    """Identify critical M&A analysis gaps"""
    print("\n" + "="*80)
    print("CRITICAL M&A ANALYSIS GAPS")
    print("="*80)
    
    critical_components = {
        'Accretion/Dilution Analysis': {
            'status': 'MISSING',
            'importance': 'CRITICAL',
            'description': 'EPS impact analysis showing whether deal is accretive or dilutive',
            'required_data': ['Acquirer financials', 'Target financials', 'Deal structure', 'Financing mix', 'Cost of capital']
        },
        'Sources & Uses': {
            'status': 'MISSING',
            'importance': 'CRITICAL',
            'description': 'Detailed breakdown of deal funding sources and uses',
            'required_data': ['Purchase price', 'Debt financing', 'Equity financing', 'Transaction fees']
        },
        'Pro Forma Financials': {
            'status': 'PARTIAL',
            'importance': 'CRITICAL',
            'description': 'Combined entity financials post-merger',
            'required_data': ['Target financials', 'Acquirer financials', 'Synergies', 'Purchase accounting']
        },
        'Contribution Analysis': {
            'status': 'MISSING',
            'importance': 'HIGH',
            'description': 'Revenue/EBITDA/EPS contribution from each party',
            'required_data': ['Standalone financials for both parties']
        },
        'Exchange Ratio Analysis': {
            'status': 'MISSING',
            'importance': 'HIGH',
            'description': 'Stock exchange ratio and premium analysis',
            'required_data': ['Stock prices', 'Shares outstanding', 'Exchange ratio']
        },
        'Financing Structure': {
            'status': 'PARTIAL',
            'importance': 'CRITICAL',
            'description': 'Debt/equity mix and financing costs',
            'required_data': ['Debt capacity', 'Interest rates', 'Equity dilution']
        },
        'Synergy Validation': {
            'status': 'PARTIAL',
            'importance': 'HIGH',
            'description': 'Detailed synergy assumptions and validation',
            'required_data': ['Cost synergies', 'Revenue synergies', 'Timeline', 'Probability']
        },
        'Fairness Opinion': {
            'status': 'MISSING',
            'importance': 'MEDIUM',
            'description': 'Fairness opinion framework and analysis',
            'required_data': ['Valuation ranges', 'Precedent analysis', 'Premiums analysis']
        }
    }
    
    for component, details in critical_components.items():
        print(f"\n{component}:")
        print(f"  Status: {details['status']}")
        print(f"  Importance: {details['importance']}")
        print(f"  Description: {details['description']}")
        print(f"  Required Data: {', '.join(details['required_data'])}")
    
    return critical_components

def generate_improvement_recommendations():
    """Generate comprehensive improvement recommendations"""
    print("\n" + "="*80)
    print("IMPROVEMENT RECOMMENDATIONS")
    print("="*80)
    
    recommendations = {
        'Phase 1 - Critical M&A Components (Immediate)': [
            '1. Implement Accretion/Dilution Analysis Agent',
            '   - Create dedicated agent for EPS impact analysis',
            '   - Calculate pro forma EPS with various financing scenarios',
            '   - Show breakeven synergy levels',
            '   - Generate sensitivity tables',
            '',
            '2. Add Sources & Uses Table Generator',
            '   - Purchase price allocation',
            '   - Debt and equity financing breakdown',
            '   - Transaction costs and fees',
            '   - Working capital adjustments',
            '',
            '3. Build Pro Forma Financial Model',
            '   - Combined income statement',
            '   - Combined balance sheet',
            '   - Combined cash flow statement',
            '   - Purchase accounting adjustments',
            '   - Goodwill calculation'
        ],
        'Phase 2 - Enhanced Analysis (High Priority)': [
            '4. Develop Contribution Analysis',
            '   - Revenue contribution by entity',
            '   - EBITDA contribution analysis',
            '   - EPS contribution breakdown',
            '   - Ownership analysis',
            '',
            '5. Create Exchange Ratio Calculator',
            '   - Calculate implied exchange ratios',
            '   - Premium/discount analysis',
            '   - Trading range analysis',
            '   - Relative valuation comparison',
            '',
            '6. Enhance Synergy Analysis',
            '   - Detailed synergy waterfall',
            '   - Probability-weighted synergies',
            '   - Synergy realization timeline',
            '   - Risk-adjusted synergy values'
        ],
        'Phase 3 - Data Flow & Integration (Medium Priority)': [
            '7. Fix Data Pipeline Issues',
            '   - Ensure all agent outputs flow to synthesis',
            '   - Add data validation checkpoints',
            '   - Implement fallback mechanisms',
            '   - Add data traceability',
            '',
            '8. Improve Report Generation',
            '   - Add more dynamic charts/visualizations',
            '   - Implement executive summary auto-generation',
            '   - Add cross-referencing between sections',
            '   - Improve formatting consistency',
            '',
            '9. Add Acquirer Company Analysis',
            '   - Currently only analyzes target',
            '   - Need parallel analysis of acquirer',
            '   - Required for pro forma and accretion/dilution'
        ],
        'Phase 4 - Advanced Features (Future)': [
            '10. Implement Fairness Opinion Framework',
            '11. Add Market Reaction Analysis',
            '12. Create Deal Scorecard',
            '13. Build Interactive Dashboard',
            '14. Add Scenario Comparison Tool'
        ]
    }
    
    for phase, items in recommendations.items():
        print(f"\n{phase}:")
        for item in items:
            print(f"  {item}")
    
    return recommendations

def main():
    """Main diagnostic function"""
    report_id = 'f02d8bcf-88ab-46c2-b516-410ee9ac9fc1'
    base_path = Path('frontend_results')
    
    print(f"Analyzing reports for job: {report_id}")
    print(f"Timestamp: {datetime.now().isoformat()}")
    
    # Analyze each report type
    excel_path = base_path / f'report_{report_id}.xlsx'
    ppt_path = base_path / f'report_{report_id}.pptx'
    pdf_path = base_path / f'report_{report_id}.pdf'
    
    results = {}
    
    if excel_path.exists():
        results['excel'] = analyze_excel_report(excel_path)
    
    if ppt_path.exists():
        results['ppt'] = analyze_ppt_report(ppt_path)
    
    if pdf_path.exists():
        results['pdf'] = analyze_pdf_report(pdf_path)
    
    # Identify gaps
    results['critical_gaps'] = identify_critical_gaps()
    
    # Generate recommendations
    results['recommendations'] = generate_improvement_recommendations()
    
    # Save results
    output_file = f'REPORT_COMPLETENESS_ANALYSIS_{report_id}.json'
    with open(output_file, 'w') as f:
        json.dump(results, f, indent=2, default=str)
    
    print(f"\n\nFull analysis saved to: {output_file}")
    
    return results

if __name__ == '__main__':
    main()
