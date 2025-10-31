"""
Detailed Placeholder and Missing Data Diagnostic
Examines actual report content to identify placeholders and data mapping issues
"""

import openpyxl
import json
from pathlib import Path
from datetime import datetime
import re

def find_excel_placeholders(filepath):
    """Find all placeholders and missing data in Excel report"""
    print("\n" + "="*80)
    print("EXCEL REPORT - PLACEHOLDER & MISSING DATA ANALYSIS")
    print("="*80)
    
    wb = openpyxl.load_workbook(filepath)
    
    issues = {
        'placeholders': [],
        'missing_data': [],
        'empty_sections': [],
        'n_a_values': []
    }
    
    # Placeholder patterns to look for
    placeholder_patterns = [
        r'\[.*?\]',  # [PLACEHOLDER]
        r'\{.*?\}',  # {PLACEHOLDER}
        r'TBD',
        r'TODO',
        r'PLACEHOLDER',
        r'XXX',
        r'N/A - .*',
        r'Data unavailable',
        r'Under analysis',
        r'Not available',
        r'Coming soon'
    ]
    
    for sheet in wb.worksheets:
        sheet_name = sheet.title
        print(f"\n{'='*60}")
        print(f"Sheet: {sheet_name}")
        print(f"{'='*60}")
        
        sheet_issues = {
            'placeholders': [],
            'missing_data': [],
            'empty_cells_in_data_area': 0,
            'n_a_count': 0
        }
        
        # Check all cells
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
            for col_idx, cell_value in enumerate(row, 1):
                if cell_value is None:
                    continue
                    
                cell_str = str(cell_value)
                
                # Check for placeholders
                for pattern in placeholder_patterns:
                    if re.search(pattern, cell_str, re.IGNORECASE):
                        sheet_issues['placeholders'].append({
                            'cell': f'{chr(64+col_idx)}{row_idx}',
                            'value': cell_str,
                            'pattern': pattern
                        })
                
                # Check for N/A
                if cell_str.upper() in ['N/A', 'NA', '#N/A']:
                    sheet_issues['n_a_count'] += 1
                
                # Check for specific missing data indicators
                if any(indicator in cell_str.lower() for indicator in [
                    'missing', 'not found', 'unavailable', 'no data',
                    'pending', 'in progress', 'under review'
                ]):
                    sheet_issues['missing_data'].append({
                        'cell': f'{chr(64+col_idx)}{row_idx}',
                        'value': cell_str
                    })
        
        # Report sheet issues
        if sheet_issues['placeholders']:
            print(f"\n🔴 PLACEHOLDERS FOUND: {len(sheet_issues['placeholders'])}")
            for item in sheet_issues['placeholders'][:5]:  # Show first 5
                print(f"   {item['cell']}: {item['value'][:50]}")
            if len(sheet_issues['placeholders']) > 5:
                print(f"   ... and {len(sheet_issues['placeholders']) - 5} more")
        
        if sheet_issues['missing_data']:
            print(f"\n⚠️  MISSING DATA INDICATORS: {len(sheet_issues['missing_data'])}")
            for item in sheet_issues['missing_data'][:5]:
                print(f"   {item['cell']}: {item['value'][:50]}")
            if len(sheet_issues['missing_data']) > 5:
                print(f"   ... and {len(sheet_issues['missing_data']) - 5} more")
        
        if sheet_issues['n_a_count'] > 0:
            print(f"\n📊 N/A VALUES: {sheet_issues['n_a_count']}")
        
        if not sheet_issues['placeholders'] and not sheet_issues['missing_data'] and sheet_issues['n_a_count'] == 0:
            print("\n✅ NO ISSUES FOUND")
        
        issues['placeholders'].extend(sheet_issues['placeholders'])
        issues['missing_data'].extend(sheet_issues['missing_data'])
        issues['n_a_values'].append({
            'sheet': sheet_name,
            'count': sheet_issues['n_a_count']
        })
    
    return issues

def find_ppt_placeholders(filepath):
    """Find placeholders in PowerPoint slides"""
    print("\n" + "="*80)
    print("POWERPOINT REPORT - PLACEHOLDER & MISSING DATA ANALYSIS")
    print("="*80)
    
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        
        issues = {
            'placeholders': [],
            'missing_data': [],
            'empty_slides': []
        }
        
        placeholder_patterns = [
            r'\[.*?\]',
            r'\{.*?\}',
            r'TBD',
            r'TODO',
            r'PLACEHOLDER',
            r'XXX',
            r'N/A - .*',
            r'Data unavailable',
            r'Under analysis',
            r'Not available'
        ]
        
        for i, slide in enumerate(prs.slides, 1):
            slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {i}"
            print(f"\n{'='*60}")
            print(f"Slide {i}: {slide_title}")
            print(f"{'='*60}")
            
            slide_issues = {
                'placeholders': [],
                'missing_data': [],
                'has_content': False
            }
            
            # Check all text in shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            if text and text.strip():
                                slide_issues['has_content'] = True
                                
                                # Check for placeholders
                                for pattern in placeholder_patterns:
                                    if re.search(pattern, text, re.IGNORECASE):
                                        slide_issues['placeholders'].append({
                                            'text': text[:100],
                                            'pattern': pattern
                                        })
                                
                                # Check for missing data indicators
                                if any(indicator in text.lower() for indicator in [
                                    'missing', 'not found', 'unavailable', 'no data',
                                    'pending', 'in progress', 'under review', 'n/a -'
                                ]):
                                    slide_issues['missing_data'].append(text[:100])
                
                # Check tables
                if shape.has_table:
                    slide_issues['has_content'] = True
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text = cell.text
                            if text:
                                for pattern in placeholder_patterns:
                                    if re.search(pattern, text, re.IGNORECASE):
                                        slide_issues['placeholders'].append({
                                            'text': text[:100],
                                            'pattern': pattern
                                        })
            
            # Report findings
            if slide_issues['placeholders']:
                print(f"\n🔴 PLACEHOLDERS: {len(slide_issues['placeholders'])}")
                for item in slide_issues['placeholders'][:3]:
                    print(f"   {item['text']}")
            
            if slide_issues['missing_data']:
                print(f"\n⚠️  MISSING DATA: {len(slide_issues['missing_data'])}")
                for item in slide_issues['missing_data'][:3]:
                    print(f"   {item}")
            
            if not slide_issues['has_content']:
                print("\n❌ EMPTY SLIDE (no content)")
                issues['empty_slides'].append(i)
            
            if not slide_issues['placeholders'] and not slide_issues['missing_data'] and slide_issues['has_content']:
                print("\n✅ NO ISSUES FOUND")
            
            issues['placeholders'].extend(slide_issues['placeholders'])
            issues['missing_data'].extend(slide_issues['missing_data'])
        
        return issues
        
    except ImportError:
        print("python-pptx not available")
        return {'error': 'python-pptx not installed'}

def find_pdf_placeholders(filepath):
    """Find placeholders in PDF"""
    print("\n" + "="*80)
    print("PDF REPORT - PLACEHOLDER & MISSING DATA ANALYSIS")
    print("="*80)
    
    try:
        import PyPDF2
        
        issues = {
            'placeholders': [],
            'missing_data': [],
            'pages_with_issues': []
        }
        
        placeholder_patterns = [
            r'\[.*?\]',
            r'\{.*?\}',
            r'TBD',
            r'TODO',
            r'PLACEHOLDER',
            r'XXX',
            r'N/A - .*',
            r'Data unavailable',
            r'Under analysis',
            r'Not available'
        ]
        
        with open(filepath, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            for i, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                
                page_issues = {
                    'placeholders': [],
                    'missing_data': []
                }
                
                # Check for placeholders
                for pattern in placeholder_patterns:
                    matches = re.findall(pattern, text, re.IGNORECASE)
                    for match in matches:
                        page_issues['placeholders'].append(match[:100])
                
                # Check for missing data indicators
                if any(indicator in text.lower() for indicator in [
                    'missing', 'not found', 'unavailable', 'no data',
                    'pending', 'in progress', 'under review', 'n/a -'
                ]):
                    page_issues['missing_data'].append(f"Page {i}")
                
                # Report findings
                if page_issues['placeholders'] or page_issues['missing_data']:
                    print(f"\nPage {i}:")
                    if page_issues['placeholders']:
                        print(f"  🔴 Placeholders: {len(page_issues['placeholders'])}")
                    if page_issues['missing_data']:
                        print(f"  ⚠️  Missing data indicators found")
                    issues['pages_with_issues'].append(i)
                
                issues['placeholders'].extend(page_issues['placeholders'])
                issues['missing_data'].extend(page_issues['missing_data'])
        
        return issues
        
    except ImportError:
        print("PyPDF2 not available")
        return {'error': 'PyPDF2 not installed'}
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return {'error': str(e)}


def analyze_synthesis_mapping():
    """Analyze synthesis mapping to identify dropped data"""
    print("\n" + "="*80)
    print("SYNTHESIS MAPPING ANALYSIS")
    print("="*80)
    
    print("""
    KNOWN SYNTHESIS MAPPING ISSUES:
    
    1. Competitive Benchmarking Data:
       - Agent outputs competitive analysis
       - Synthesis sometimes fails to extract it
       - Location: state['competitive_benchmarking'] or agent_outputs array
       - Fix needed in: synthesis_reporting.py/_find_agent_data()
    
    2. Macroeconomic Analysis Data:
       - Agent outputs scenario models
       - Synthesis may miss correlation analysis
       - Location: state['macroeconomic_analyst']
       - Fix needed in: synthesis_reporting.py/_generate_risk_macro_section()
    
    3. Financial Deep Dive Ratios:
       - Agent calculates detailed ratios
       - Synthesis may not fully extract all ratios
       - Location: state['financial_deep_dive']
       - Fix needed in: synthesis_reporting.py/_generate_financial_section()
    
    4. Integration Planning Synergies:
       - Agent creates detailed synergy breakdown
       - Synthesis may use only summary
       - Location: state['integration_planner']
       - Fix needed in: synthesis_reporting.py/_generate_integration_tax_section()
    
    5. Tax Structuring Details:
       - Agent provides comprehensive tax analysis
       - Synthesis may miss optimization opportunities
       - Location: state['tax_structuring']
       - Fix needed in: synthesis_reporting.py/_generate_integration_tax_section()
    
    6. External Validation Results:
       - Agent performs external cross-checks
       - Synthesis needs dedicated extraction method
       - Location: state['external_validator']
       - Fix: Added _generate_external_validation_section() TODAY
    """)


def generate_fix_plan():
    """Generate detailed fix plan for placeholders and missing data"""
    print("\n" + "="*80)
    print("DETAILED FIX PLAN FOR PLACEHOLDERS & MISSING DATA")
    print("="*80)
    
    fixes = {
        'Excel Report Fixes': [
            {
                'issue': 'N/A values in Competitive Benchmarking sheet',
                'root_cause': 'Synthesis not extracting competitive_benchmarking agent data',
                'file_to_fix': 'src/agents/synthesis_reporting.py',
                'function': '_generate_market_section()',
                'fix': 'Enhance data extraction from state[\'competitive_benchmarking\']',
                'priority': 'HIGH'
            },
            {
                'issue': 'Placeholder text "Under analysis" in various sheets',
                'root_cause': 'Fallback data being used when agent data not found',
                'file_to_fix': 'src/agents/synthesis_reporting.py',
                'function': '_find_agent_data() and _collect_fallback_agent_data()',
                'fix': 'Improve agent data search and remove fallback generation',
                'priority': 'CRITICAL'
            },
            {
                'issue': 'Missing macro scenario data',
                'root_cause': 'macroeconomic_analyst data not fully extracted',
                'file_to_fix': 'src/agents/synthesis_reporting.py',
                'function': '_generate_risk_macro_section()',
                'fix': 'Extract scenario_models, correlation_analysis from macro agent',
                'priority': 'HIGH'
            },
            {
                'issue': 'Empty cells in financial ratios',
                'root_cause': 'financial_deep_dive ratios not all mapped',
                'file_to_fix': 'src/outputs/excel_generator.py',
                'function': '_build_financial_deep_dive_sheet()',
                'fix': 'Map all ratio categories from synthesis data',
                'priority': 'MEDIUM'
            }
        ],
        'PPT Report Fixes': [
            {
                'issue': '"N/A - Data unavailable" on competitive slides',
                'root_cause': 'Same as Excel - synthesis mapping issue',
                'file_to_fix': 'src/outputs/ppt_sections/financial_slides.py',
                'function': 'add_competitive_benchmarking_slide()',
                'fix': 'Check synthesis.market_analysis.competitive_landscape',
                'priority': 'HIGH'
            },
            {
                'issue': 'Generic placeholder text in macro slides',
                'root_cause': 'Macro data not fully populated',
                'file_to_fix': 'src/outputs/ppt_sections/risk_slides.py',
                'function': 'add_macroeconomic_slide()',
                'fix': 'Extract scenario data from synthesis.risk_macro',
                'priority': 'HIGH'
            },
            {
                'issue': 'Missing synergy details',
                'root_cause': 'Integration synergy breakdown not extracted',
                'file_to_fix': 'src/outputs/ppt_sections/executive_slides.py',
                'function': 'add_synergy_slide()',
                'fix': 'Use synthesis.integration_tax.synergy_breakdown',
                'priority': 'MEDIUM'
            }
        ],
        'PDF Report Fixes': [
            {
                'issue': 'Sections with minimal content',
                'root_cause': 'Same synthesis issues as Excel/PPT',
                'file_to_fix': 'src/outputs/pdf_sections/*',
                'function': 'Multiple section generators',
                'fix': 'Ensure all sections pull from synthesis correctly',
                'priority': 'HIGH'
            },
            {
                'issue': 'Tables with N/A values',
                'root_cause': 'Data not flowing from synthesis',
                'file_to_fix': 'src/outputs/revolutionary_pdf_generator.py',
                'function': '_generate_*_section() methods',
                'fix': 'Add fallback to [] or "Not assessed" instead of "N/A"',
                'priority': 'MEDIUM'
            }
        ],
        'Synthesis Agent Fixes': [
            {
                'issue': 'Agent data not being found',
                'root_cause': '_find_agent_data() not searching all locations',
                'file_to_fix': 'src/agents/synthesis_reporting.py',
                'function': '_find_agent_data()',
                'fix': 'Enhanced recursive search - PARTIALLY DONE TODAY',
                'priority': 'CRITICAL',
                'status': 'IN PROGRESS'
            },
            {
                'issue': 'Fallback data being used too often',
                'root_cause': 'Agent data exists but not being found',
                'file_to_fix': 'src/agents/synthesis_reporting.py',
                'function': '_collect_fallback_agent_data()',
                'fix': 'Remove fallback generation, throw error instead',
                'priority': 'HIGH'
            }
        ]
    }
    
    for category, issue_list in fixes.items():
        print(f"\n{category}:")
        print("="*60)
        for issue in issue_list:
            print(f"\n  Issue: {issue['issue']}")
            print(f"  Root Cause: {issue['root_cause']}")
            print(f"  File: {issue['file_to_fix']}")
            print(f"  Function: {issue['function']}")
            print(f"  Fix: {issue['fix']}")
            print(f"  Priority: {issue['priority']}")
            if 'status' in issue:
                print(f"  Status: {issue['status']}")
    
    return fixes


def main():
    """Main diagnostic function"""
    report_id = 'f02d8bcf-88ab-46c2-b516-410ee9ac9fc1'
    base_path = Path('frontend_results')
    
    print(f"PLACEHOLDER & MISSING DATA DIAGNOSTIC")
    print(f"Report ID: {report_id}")
    print(f"Timestamp: {datetime.now().isoformat()}")
    print("="*80)
    
    results = {
        'report_id': report_id,
        'timestamp': datetime.now().isoformat()
    }
    
    # Analyze Excel
    excel_path = base_path / f'report_{report_id}.xlsx'
    if excel_path.exists():
        results['excel'] = find_excel_placeholders(excel_path)
    
    # Analyze PPT
    ppt_path = base_path / f'report_{report_id}.pptx'
    if ppt_path.exists():
        results['ppt'] = find_ppt_placeholders(ppt_path)
    
    # Analyze PDF
    pdf_path = base_path / f'report_{report_id}.pdf'
    if pdf_path.exists():
        results['pdf'] = find_pdf_placeholders(pdf_path)
    
    # Analyze synthesis mapping
    analyze_synthesis_mapping()
    
    # Generate fix plan
    results['fix_plan'] = generate_fix_plan()
    
    # Save results
    output_file = f'PLACEHOLDER_ANALYSIS_{report_id}.json'
    with open(output_file, 'w') as f:
        json.dump(results, f, indent=2, default=str)
    
    print(f"\n\n{'='*80}")
    print(f"Full analysis saved to: {output_file}")
    print("="*80)
    
    # Summary
    excel_issues = results.get('excel', {})
    ppt_issues = results.get('ppt', {})
    pdf_issues = results.get('pdf', {})
    
    print("\n\nSUMMARY:")
    print("="*80)
    print(f"Excel Placeholders: {len(excel_issues.get('placeholders', []))}")
    print(f"Excel Missing Data: {len(excel_issues.get('missing_data', []))}")
    print(f"PPT Placeholders: {len(ppt_issues.get('placeholders', []))}")
    print(f"PPT Missing Data: {len(ppt_issues.get('missing_data', []))}")
    print(f"PDF Placeholders: {len(pdf_issues.get('placeholders', []))}")
    print(f"PDF Pages with Issues: {len(pdf_issues.get('pages_with_issues', []))}")
    
    total_issues = (
        len(excel_issues.get('placeholders', [])) +
        len(excel_issues.get('missing_data', [])) +
        len(ppt_issues.get('placeholders', [])) +
        len(ppt_issues.get('missing_data', [])) +
        len(pdf_issues.get('placeholders', []))
    )
    
    print(f"\n🔴 TOTAL ISSUES IDENTIFIED: {total_issues}")
    print("\nRecommendation: Review fix plan above and prioritize CRITICAL items")
    
    return results


if __name__ == '__main__':
    main()
